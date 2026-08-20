import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

from config import OUTPUT_DIR
from core.state import DeckState, SlideState, SlideElement
from core.logger import logger

def hex_to_rgb(hex_str: str) -> RGBColor:
    """Converts hex color code e.g. '#0D6EFD' to pptx RGBColor."""
    hex_str = hex_str.lstrip('#')
    if len(hex_str) == 6:
        r, g, b = tuple(int(hex_str[i:i+2], 16) for i in (0, 2, 4))
        return RGBColor(r, g, b)
    return RGBColor(33, 37, 41)

def parse_dimension(dim_str: str, default_val: float, total_span: float) -> Inches:
    """Converts dimension string like '5%', '0.8in', '40pt' to pptx Inches."""
    try:
        dim_str = dim_str.strip().lower()
        if dim_str.endswith('%'):
            pct = float(dim_str.rstrip('%')) / 100.0
            return Inches(pct * total_span)
        elif dim_str.endswith('in'):
            return Inches(float(dim_str.rstrip('in')))
        elif dim_str.endswith('pt'):
            return Inches(float(dim_str.rstrip('pt')) / 72.0)
        elif dim_str.endswith('px'):
            return Inches(float(dim_str.rstrip('px')) / 96.0)
    except Exception:
        pass
    return Inches(default_val)

def build_pptx_from_deck_state(deck_state: DeckState, filename: str = "final_presentation.pptx") -> str:
    """
    Builds a native, fully editable PowerPoint presentation (.pptx) from the DeckState JSON specification.
    Every element (headline, paragraph, bullet point, diagram, photo) is created as an individual editable shape.
    """
    prs = Presentation()
    
    # Set 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Theme settings
    theme = deck_state.theme
    bg_color = hex_to_rgb(theme.background_color)
    text_color = hex_to_rgb(theme.text_color)
    primary_color = hex_to_rgb(theme.primary_color)
    
    blank_layout = prs.slide_layouts[6] # Blank slide layout
    
    for slide_data in deck_state.slides:
        slide = prs.slides.add_slide(blank_layout)
        
        # Apply background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        # Track created text frames by position to prevent stacking overlapping textboxes
        created_text_frames = {} # (round_left, round_top) -> tf
        
        for element in slide_data.elements:
            left = parse_dimension(element.position.left, 0.8, 13.333)
            top = parse_dimension(element.position.top, 0.8, 7.5)
            width = parse_dimension(element.position.width, 11.7, 13.333)
            height = parse_dimension(element.position.height, 5.0, 7.5)
            
            if element.type == "text":
                has_bg = bool(element.style and element.style.background_color)
                
                # Check if a text box already exists at roughly this same position
                existing_tf = None
                if not has_bg: # Only merge standard text boxes, keep distinct cards separate
                    for (ex_left, ex_top), ex_tf in created_text_frames.items():
                        if abs(ex_left - left) < Inches(0.3) and abs(ex_top - top) < Inches(0.3):
                            existing_tf = ex_tf
                            break
                
                if existing_tf is not None:
                    tf = existing_tf
                    is_new_box = False
                else:
                    # Create new textbox shape
                    if has_bg:
                        txBox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
                        txBox.fill.solid()
                        try:
                            txBox.fill.fore_color.rgb = hex_to_rgb(element.style.background_color)
                            txBox.line.fill.background()
                        except:
                            pass
                    else:
                        txBox = slide.shapes.add_textbox(left, top, width, height)
                    
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    tf.margin_left = Inches(0.2) if has_bg else Inches(0)
                    tf.margin_right = Inches(0.2) if has_bg else Inches(0)
                    tf.margin_top = Inches(0.2) if has_bg else Inches(0)
                    tf.margin_bottom = Inches(0.2) if has_bg else Inches(0)
                    
                    if not has_bg:
                        created_text_frames[(left, top)] = tf
                    is_new_box = True
                
                content = element.content
                if isinstance(content, str):
                    lines = [content]
                elif isinstance(content, list):
                    lines = content
                else:
                    lines = []
                    
                for i, line in enumerate(lines):
                    if is_new_box and i == 0 and len(tf.paragraphs) > 0 and tf.paragraphs[0].text == "":
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                        
                    p.text = str(line)
                    
                    # --- DEFAULT ROLE STYLING ---
                    default_font_name = theme.font_body
                    default_font_size = 18
                    default_font_bold = False
                    default_font_color = text_color
                    
                    if element.role == "slide_title":
                        default_font_size = 32
                        default_font_bold = True
                        default_font_color = primary_color
                        default_font_name = theme.font_title
                    elif element.role == "subtitle":
                        default_font_size = 24
                        default_font_bold = True
                        default_font_color = primary_color
                        default_font_name = theme.font_title
                    elif element.role == "bullet_list":
                        default_font_size = 18
                        default_font_bold = False
                        default_font_color = text_color
                        p.level = 0
                        
                    # --- APPLY STYLES ---
                    # WE STRICLY USE DEFAULT FONT FAMILIES FOR CONSISTENCY
                    p.font.name = default_font_name
                    
                    # Parse dynamic font size from LLM if provided, otherwise default
                    parsed_size = default_font_size
                    if element.style and element.style.font_size:
                        size_str = str(element.style.font_size).lower().replace('pt', '').replace('px', '').strip()
                        try:
                            parsed_size = int(float(size_str))
                        except Exception:
                            pass
                    
                    # Clamp font sizes to safe, non-overlapping bounds
                    if element.role == "slide_title":
                        parsed_size = min(max(parsed_size, 28), 36)
                    elif element.role == "subtitle":
                        parsed_size = min(max(parsed_size, 20), 24)
                    elif element.role in ["bullet_list", "body_text"]:
                        parsed_size = min(max(parsed_size, 14), 18)
                    
                    p.font.size = Pt(parsed_size)
                    
                    # Force PowerPoint to autofit the text if it overflows the bounds!
                    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                            
                    # Font Weight
                    if element.style and element.style.font_weight in ["bold", "700", "800", "900"]:
                        p.font.bold = True
                    else:
                        p.font.bold = default_font_bold
                        
                    # Font Color
                    if element.style and element.style.color:
                        try:
                            p.font.color.rgb = hex_to_rgb(element.style.color)
                        except:
                            p.font.color.rgb = default_font_color
                    else:
                        p.font.color.rgb = default_font_color
                        
                    # Alignment
                    if element.style and element.style.alignment:
                        align_str = element.style.alignment.lower()
                        if align_str == "center":
                            p.alignment = PP_ALIGN.CENTER
                        elif align_str == "right":
                            p.alignment = PP_ALIGN.RIGHT
                        elif align_str == "justify":
                            p.alignment = PP_ALIGN.JUSTIFY
                        else:
                            p.alignment = PP_ALIGN.LEFT
                        
            elif element.type and ("image" in element.type.lower() or "diagram" in element.type.lower()):
                if element.file_path and os.path.exists(element.file_path):
                    try:
                        slide.shapes.add_picture(element.file_path, left, top, width, height)
                        logger.info(f"Added picture shape ({element.type}) at slide {slide_data.slide_number}")
                    except Exception as e:
                        logger.error(f"Failed to place image shape: {e}")
                        
        # Render programmatic footer for non-title slides
        if slide_data.layout_type not in ["title_slide", "section_divider", "closing_slide"]:
            # Slide Number
            num_left, num_top = Inches(12.5), Inches(7.0)
            txBox = slide.shapes.add_textbox(num_left, num_top, Inches(0.5), Inches(0.3))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = str(slide_data.slide_number)
            p.font.size = Pt(12)
            p.font.name = theme.font_body
            p.font.color.rgb = text_color
            p.alignment = PP_ALIGN.RIGHT
            
            # Deck Title / Footer Text
            footer_text = theme.footer_text or deck_state.user_prompt[:50] + "..."
            title_left, title_top = Inches(0.5), Inches(7.0)
            txBoxTitle = slide.shapes.add_textbox(title_left, title_top, Inches(5.0), Inches(0.3))
            tfTitle = txBoxTitle.text_frame
            pTitle = tfTitle.paragraphs[0]
            pTitle.text = footer_text
            pTitle.font.size = Pt(12)
            pTitle.font.name = theme.font_body
            pTitle.font.color.rgb = hex_to_rgb(theme.accent_color)
            pTitle.alignment = PP_ALIGN.LEFT
            
    output_path = os.path.join(OUTPUT_DIR, filename)
    prs.save(output_path)
    logger.info(f"Successfully generated native editable PPTX file -> {output_path}")
    return output_path
